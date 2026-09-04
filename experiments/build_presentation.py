"""Generate an editable PowerPoint (.pptx) version of the project
presentation -- real content/narrative, in a format you can open in
PowerPoint/Keynote/Google Slides and edit directly.

Flow: problem statement -> the technique that already exists (JA3/JA3S)
-> what we were asked to build -> architecture -> database -> results
-> an unexpected discovery -> JA4 fix -> bot attacks -> spoofing
detection -> stress test -> limitations -> conclusion.

Uses only fonts bundled with Microsoft Office (Calibri / Consolas) so it
renders correctly in actual PowerPoint/Keynote, not just in a browser.

Usage:
    python experiments/build_presentation.py
    # writes TLS_Fingerprinting_Presentation.pptx to the project root
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "TLS_Fingerprinting_Presentation.pptx"

# ---- palette ----------------------------------------------------------
BG = RGBColor(0x0A, 0x0E, 0x17)
CARD = RGBColor(0x12, 0x18, 0x26)
BORDER = RGBColor(0x23, 0x2B, 0x3D)
TEXT = RGBColor(0xED, 0xEF, 0xF4)
MUTED = RGBColor(0x9A, 0xA3, 0xB5)
FAINT = RGBColor(0x69, 0x72, 0x8A)
ACCENT = RGBColor(0xE8, 0xA3, 0x3D)
OK = RGBColor(0x49, 0xC7, 0x9A)
FLAG = RGBColor(0xF1, 0x6A, 0x54)

# Office-bundled fonts only -- renders correctly outside the browser.
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs, number, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    num_box = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.55), Inches(2), Inches(0.4))
    tf = num_box.text_frame
    tf.text = f"{number:02d} / {total}"
    run = tf.paragraphs[0].runs[0]
    run.font.name = FONT_MONO
    run.font.size = Pt(11)
    run.font.color.rgb = FAINT
    return slide


def add_textbox(slide, left, top, width, height, text, font=FONT_BODY,
                 size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
                 line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
    return box


def add_eyebrow(slide, text, top=Inches(0.65)):
    add_textbox(slide, MARGIN, top, Inches(11), Inches(0.4), text.upper(),
                font=FONT_MONO, size=13, color=ACCENT, bold=True)


def add_title(slide, text, top=Inches(1.15), size=34, width=Inches(11.8)):
    add_textbox(slide, MARGIN, top, width, Inches(1.7), text,
                font=FONT_TITLE, size=size, color=TEXT, bold=True, line_spacing=1.05)


def add_body(slide, text, top, size=16, width=Inches(10.5), color=MUTED):
    add_textbox(slide, MARGIN, top, width, Inches(2), text,
                font=FONT_BODY, size=size, color=color, line_spacing=1.35)


def add_bullets(slide, items, top, size=17, width=Inches(10.8), gap=Inches(0.62)):
    y = top
    for text in items:
        box = slide.shapes.add_textbox(MARGIN, y, width, gap)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"•  {text}"
        p.line_spacing = 1.2
        for run in p.runs:
            run.font.name = FONT_BODY
            run.font.size = Pt(size)
            run.font.color.rgb = MUTED
        y = y + gap
    return y


def add_card(slide, left, top, width, height, lines, title=None):
    """lines: list of (text, font, size, color, bold)"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.28)
    tf.margin_right = Inches(0.28)
    tf.margin_top = Inches(0.22)
    tf.margin_bottom = Inches(0.22)

    all_lines = []
    if title:
        all_lines.append((title, FONT_MONO, 12, FAINT, False))
    all_lines.extend(lines)

    for i, (text, font, size, color, bold) in enumerate(all_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.line_spacing = 1.25
        p.space_after = Pt(5)
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
    return card


def add_kpi_row(slide, items, top):
    left = MARGIN
    w = Inches(2.75)
    for n, label in items:
        add_textbox(slide, left, top, w, Inches(0.7), n, font=FONT_TITLE,
                    size=32, color=ACCENT, bold=True)
        add_textbox(slide, left, top + Inches(0.65), w, Inches(0.5), label,
                    font=FONT_BODY, size=12, color=MUTED)
        left = left + w


def add_table(slide, left, top, width, height, header, rows):
    n_rows = len(rows) + 1
    n_cols = len(header)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for c, h in enumerate(header):
        cell = gtable.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD
        p = cell.text_frame.paragraphs[0]
        for run in p.runs:
            run.font.name = FONT_MONO
            run.font.size = Pt(11)
            run.font.color.rgb = FAINT
            run.font.bold = False
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.name = FONT_MONO if c > 0 else FONT_BODY
                run.font.size = Pt(12)
                run.font.color.rgb = TEXT if c == 0 else MUTED
    return gtable


def build():
    prs = new_deck()
    total = 15

    # 1. Title
    s = add_slide(prs, 1, total)
    add_eyebrow(s, "Computer Networks — Course Project", top=Inches(2.2))
    add_title(s, "TLS Fingerprinting", top=Inches(2.7), size=54)
    add_body(s, "Identifying who's really on the other end of an HTTPS "
                "connection — without decrypting anything.",
             top=Inches(3.75), size=18, width=Inches(9.5))
    add_kpi_row(s, [("5", "real clients captured"), ("15", "database entries"),
                    ("0", "fabricated results")], top=Inches(4.75))

    # 2. Problem statement
    s = add_slide(prs, 2, total)
    add_eyebrow(s, "Problem Statement")
    add_title(s, "Once HTTPS starts, a server can't read what's inside.", size=30)
    add_bullets(s, [
        "Almost all internet traffic today is encrypted (HTTPS/TLS).",
        "A bot pretending to be a browser, a scraper, or malware talking "
        "to its C2 server — all look like ordinary encrypted traffic.",
        "Firewalls and monitors can no longer read URLs, headers, or payloads.",
        "Question: can we tell WHO is really connecting, without breaking "
        "encryption for anyone?",
    ], top=Inches(2.5), size=18, gap=Inches(0.78))

    # 3. What already exists — JA3/JA3S
    s = add_slide(prs, 3, total)
    add_eyebrow(s, "What Already Exists")
    add_title(s, "JA3 / JA3S: a published fingerprinting technique", size=30)
    add_bullets(s, [
        'Right before encryption locks in, both sides send ONE message '
        'each in plain text: the "ClientHello" and "ServerHello".',
        "Every TLS library (Chrome, curl, OpenSSL...) builds this "
        "handshake slightly differently — different ciphers, different order.",
        "JA3 (Salesforce, 2017) turns a ClientHello into one short hash. "
        "JA3S does the same for a ServerHello.",
        "This is a real, industry-standard technique — not something we "
        "invented for this project.",
    ], top=Inches(2.5), size=17, gap=Inches(0.78))

    # 4. What we were asked to build
    s = add_slide(prs, 4, total)
    add_eyebrow(s, "Project Scope")
    add_title(s, "What we were asked to implement", size=32)
    add_bullets(s, [
        "Passively capture TLS handshakes (ClientHello / ServerHello)",
        "Compute JA3 (client) and JA3S (server) fingerprints",
        "Build a small reference database of known fingerprints",
        "Match a new capture against that database and identify the client",
        "Demonstrate on at least 5 real, distinct clients",
        "Explain the security use case AND the honest limitations",
    ], top=Inches(2.3), size=18, gap=Inches(0.68))

    # 5. Architecture
    s = add_slide(prs, 5, total)
    add_eyebrow(s, "How It's Built")
    add_title(s, "The pipeline", size=34)
    steps = [".pcap", "parser.py", "JA3 / JA3S / JA4", "database.py", "CLI output"]
    left = MARGIN
    top = Inches(3.1)
    w = Inches(2.05)
    for i, step in enumerate(steps):
        add_card(s, left, top, w, Inches(0.9), [(step, FONT_MONO, 13, TEXT, False)])
        left = left + w
        if i < len(steps) - 1:
            add_textbox(s, left, top + Inches(0.22), Inches(0.25), Inches(0.5), "->",
                        font=FONT_MONO, size=16, color=FAINT)
            left = left + Inches(0.25)
    add_body(s, "Every stage is a separate, independently-tested Python module — "
                "the handshake is parsed byte-by-byte against the real TLS spec, "
                "not guessed.", top=Inches(4.5), size=16, width=Inches(10.5))

    # 6. Testing & validation
    s = add_slide(prs, 6, total)
    add_eyebrow(s, "Before Trusting Any Result")
    add_title(s, "56 automated tests — what they actually check", size=28)
    add_bullets(s, [
        "JA3 / JA3S string construction, checked against hand-derived "
        "expected values (not just checked against itself).",
        "JA4, checked against the official FoxIO spec's own published "
        "worked examples — an external ground truth.",
        "Parser edge cases: truncated records, GREASE-only cipher lists, "
        "a handshake split across multiple TLS records.",
        "Database matching logic (known / possible / unknown) and the "
        "spoofing-detector's mismatch logic.",
        "One end-to-end test: build a synthetic pcap, run the full "
        "pipeline, check the final answer.",
    ], top=Inches(2.35), size=16, gap=Inches(0.68))
    add_body(s, "Runs offline in under a second, no network needed — this is "
                "what's checked BEFORE any real capture is trusted.",
             top=Inches(6.35), size=13, color=FAINT)

    # 7. The reference database
    s = add_slide(prs, 7, total)
    add_eyebrow(s, "Reference Database")
    add_title(s, "A notebook of known fingerprints", size=32)
    add_body(s, "Every entry was computed from a real, captured network "
                "connection — never typed by hand or copied from a public list.",
             top=Inches(2.35), size=17, width=Inches(10.8))
    add_card(s, MARGIN, Inches(3.2), Inches(11.5), Inches(2.9), [
        ("[ja3]  curl 8.7.1 (macOS system, SecureTransport/LibreSSL)", FONT_MONO, 13, TEXT, False),
        ("[ja4]  curl 8.7.1 (macOS system, SecureTransport/LibreSSL)", FONT_MONO, 13, TEXT, False),
        ("[ja3s] Cloudflare edge (fronting example.com)", FONT_MONO, 13, MUTED, False),
        ("[ja3]  Google Chrome 151.0.7922.174 (headless)", FONT_MONO, 13, TEXT, False),
        ("[ja4]  Google Chrome 151.0.7922.174 (headless)", FONT_MONO, 13, TEXT, False),
        ("... 15 entries total: 5 clients x (JA3 + JA4 + JA3S)", FONT_MONO, 13, FAINT, False),
    ], title="tls-fingerprint db list")

    # 8. Results — 5 clients table
    s = add_slide(prs, 8, total)
    add_eyebrow(s, "Results")
    add_title(s, "Five different clients. Five different fingerprints.", size=27)
    add_table(
        s, MARGIN, Inches(2.5), Inches(11.5), Inches(3.1),
        ["Client", "TLS library", "JA3 hash"],
        [
            ["curl 8.7.1", "SecureTransport / LibreSSL", "375c6162a492...ce8424"],
            ["openssl s_client", "OpenSSL 3.6.2", "0b85eb0d4981...f0ac5f"],
            ["Python stdlib ssl", "OpenSSL 3.6.2 (same lib!)", "f21f8e6cf70d...ef401c"],
            ["Chrome 151 (headless)", "BoringSSL", "81a2542af844...f2a626"],
            ["Hand-built ClientHello", "none — raw socket", "c53113116bb0...6fedc9"],
        ],
    )
    add_body(s, "openssl and Python use the exact same crypto library and still "
                "get different fingerprints — JA3 reflects configuration, not "
                "just which library is linked.", top=Inches(6.0), size=13, color=FAINT)

    # 9. Unexpected discovery
    s = add_slide(prs, 9, total)
    add_eyebrow(s, "An Unplanned Discovery")
    add_title(s, "The same Chrome, twice, gave two different answers.", size=28)
    add_card(s, MARGIN, Inches(2.6), Inches(5.6), Inches(1.6),
             [("81a2542af8442fcd7802f178d9f2a626", FONT_MONO, 14, TEXT, False)],
             title="Run 1")
    add_card(s, Inches(6.65), Inches(2.6), Inches(5.6), Inches(1.6),
             [("a00e551d2f4af85ede1156537ebf095a", FONT_MONO, 14, FLAG, True)],
             title="Run 2 — same install, moments later")
    add_body(s, "Modern Chrome deliberately randomizes ClientHello extension "
                "order per connection, specifically to weaken fingerprinting "
                "like this. We didn't just read that fact — we reproduced it "
                "live, re-running our own experiment.",
             top=Inches(4.6), size=16, width=Inches(10.8))

    # 10. JA4 fix + comparison
    s = add_slide(prs, 10, total)
    add_eyebrow(s, "Built In Response — JA4")
    add_title(s, "JA4: a newer fingerprint that isolates what actually changed", size=27)
    add_card(s, MARGIN, Inches(2.6), Inches(5.6), Inches(2.9), [
        ("8daaf6152771", FONT_MONO, 15, OK, True),
        ("8daaf6152771", FONT_MONO, 15, OK, True),
        ("Identical, both runs — JA4 sorts the\ncipher list before hashing, so pure\nreordering can't change it.", FONT_BODY, 12, FAINT, False),
    ], title="Cipher segment (JA4)")
    add_card(s, Inches(6.65), Inches(2.6), Inches(5.6), Inches(2.9), [
        ("806a8c22fdea   (16 extensions)", FONT_MONO, 13, TEXT, False),
        ("541cd5a3d78e   (17 extensions)", FONT_MONO, 13, FLAG, True),
        ("Genuinely different — Chrome sent one\nextra extension the 2nd time. JA4 reports\nthat honestly instead of hiding it.", FONT_BODY, 12, FAINT, False),
    ], title="Extension segment (JA4)")
    add_body(s, "Validated against the official FoxIO spec's own published "
                "worked examples — not just checked against itself.",
             top=Inches(5.75), size=13, color=FAINT)

    # 11. Bot attacks — how fake clients lie
    s = add_slide(prs, 11, total)
    add_eyebrow(s, "From Fingerprint To Defense")
    add_title(s, "How a bot pretends to be Chrome", size=32)
    add_bullets(s, [
        'The cheap, common trick: set the HTTP "User-Agent" header to '
        'Chrome\'s exact string — one line of code, free to fake.',
        "But the TLS ClientHello was already sent, moments earlier, by "
        "whichever real library the bot's script actually uses.",
        "That handshake is a structural byproduct of the real library — "
        "not something typed by hand, much harder to fake convincingly.",
        "The mismatch between the claim and the real handshake is the "
        "tell.",
    ], top=Inches(2.4), size=18, gap=Inches(0.78))

    # 12. Spoofing detection demo
    s = add_slide(prs, 12, total)
    add_eyebrow(s, "Catching The Lie")
    add_title(s, "Spoofing detection, live", size=34)
    add_card(s, MARGIN, Inches(2.5), Inches(11.5), Inches(2.6), [
        ("$ tls-fingerprint check-spoofing pcaps/bot_client.pcap --claims Chrome", FONT_MONO, 13, FAINT, False),
        ("Claims to be:  Chrome", FONT_MONO, 13, MUTED, False),
        ("Measured JA3:  f21f8e6cf70d5980ecfe9fa2e0ef401c", FONT_MONO, 13, MUTED, False),
        ("VERDICT: *** MISMATCH — SUSPICIOUS ***", FONT_MONO, 15, FLAG, True),
        ("matches: Python 3.14.6 stdlib ssl (ssl.create_default_context)", FONT_MONO, 13, MUTED, False),
    ])
    add_body(s, "The claimed identity was checked against the database's real "
                "Chrome hashes from slide 7 — it didn't match, so it's flagged, "
                "and the tool reports what it actually is instead.",
             top=Inches(5.6), size=15, width=Inches(10.8))

    # 13. Bombardment / stress test
    s = add_slide(prs, 13, total)
    add_eyebrow(s, "Under Load")
    add_title(s, "5 rapid requests. 5 for 5, caught.", size=32)
    add_card(s, MARGIN, Inches(2.4), Inches(11.5), Inches(3.0), [
        ("$ python experiments/bombard_demo.py 5", FONT_MONO, 13, FAINT, False),
        ("request 1/5: JA3=f21f8e6c...  -> FLAGGED", FONT_MONO, 13, FLAG, False),
        ("request 2/5: JA3=f21f8e6c...  -> FLAGGED", FONT_MONO, 13, FLAG, False),
        ("request 3/5: JA3=f21f8e6c...  -> FLAGGED", FONT_MONO, 13, FLAG, False),
        ("request 4/5: JA3=f21f8e6c...  -> FLAGGED", FONT_MONO, 13, FLAG, False),
        ("request 5/5: JA3=f21f8e6c...  -> FLAGGED", FONT_MONO, 13, FLAG, False),
        ("Result: 5/5 requests correctly flagged.", FONT_MONO, 14, OK, True),
    ])
    add_body(s, "Real botnets spread requests across many IPs to stay under "
                "rate-limit thresholds. That trick doesn't help here — every "
                "connection is judged on its own real handshake, independently, "
                "so more requests just means more identical evidence.",
             top=Inches(5.7), size=14, width=Inches(10.8))

    # 14. Limitations
    s = add_slide(prs, 14, total)
    add_eyebrow(s, "What This Doesn't Prove")
    add_title(s, "A fingerprint is a hint, not a verdict.", size=32)
    limits = [
        "Same hash != same software — two unrelated programs on the same "
        "library/config fingerprint identically.",
        "JA3S depends on the client, too — two of our own clients produced "
        "an identical JA3S against the same server.",
        "GREASE and reordering actively fight back — Chrome randomizes "
        "its own hello on purpose (slide 9).",
        "Evasion is possible — since the fingerprint is client-controlled "
        "bytes, a determined attacker can copy it exactly.",
        "Nothing here decrypts anything — only the already-unencrypted "
        "handshake preamble is ever read.",
    ]
    top = Inches(2.4)
    for i, text in enumerate(limits, start=1):
        add_textbox(s, MARGIN, top, Inches(0.6), Inches(0.5), f"{i:02d}",
                    font=FONT_MONO, size=14, color=ACCENT, bold=True)
        add_textbox(s, Inches(1.5), top, Inches(10.5), Inches(0.75), text,
                    font=FONT_BODY, size=15, color=MUTED, line_spacing=1.2)
        top = top + Inches(0.85)

    # 15. Conclusion
    s = add_slide(prs, 15, total)
    add_eyebrow(s, "In Summary")
    add_title(s, "Built, tested, and proven against real traffic.", size=32)
    add_body(s, "Every hash on every slide came from a real handshake with a "
                "real server — none of it invented.",
             top=Inches(2.3), size=18, width=Inches(10.5))
    add_kpi_row(s, [("15", "measured DB entries"), ("7", "real experiments"),
                    ("2", "fingerprint schemes")],
                top=Inches(3.4))
    add_body(s, "TLS Fingerprinting — Computer Networks course project.",
             top=Inches(5.2), size=13, color=FAINT)

    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    build()
